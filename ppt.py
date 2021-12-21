# -*- coding: utf-8 -*-

import pptx
from pptx import Presentation

prs = Presentation("统计学习方法PPT.pptx")
for slide in prs.slides:
	print(slide)
	for shape in slide.shapes:
		print(shape)
		if shape.has_text_frame:
			text_frame = shape.text_frame
			print(text_frame.text)

# 获取某一页Slide中的内容
for i,slide in enumerate(prs.slides):
	if i == 5:
		for shape in slide.shapes:
			if shape.has_text_frame:
				text_frame = shape.text_frame
			print(text_frame.text)
# 获取Shape中的某个Paragraph
import pptx
from pptx import Presentation

prs = Presentation("统计学习方法PPT.pptx")

for i,slide in enumerate(prs.slides):
	if i == 5:
		for shape in slide.shapes:
			if shape.has_text_frame:
				text_frame = shape.text_frame
				for paragraph in text_frame.paragraphs:
					print(paragraph.text)

# 添加Slide和内容

# 占位符id的确认
import pptx
from pptx import Presentation

prs = Presentation("空白.pptx")
# prs.slide_layouts[]表示的是ppt中不同的版式
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
	phf = shape.placeholder_format
	print(f"{phf.idx}--{shape.name}--{phf.type}")
	shape.text = f"{phf.idx}--{shape.name}--{phf.type}"
	# 注意：做完这个操作，一定要记得保存一下！
	prs.save("电子奖状模板.pptx")

# PPT内容的填写
import pptx
from pptx import Presentation

prs = Presentation("空白.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
name = slide.placeholders[14]
why = slide.placeholders[15]

name.text = "黄同学"
why.text = "学习太积极"
prs.save("内容填充.pptx")

# 添加段落
# 占位符id的确认
import pptx
from pptx import Presentation

prs = Presentation("finall.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
	phf = shape.placeholder_format
	print(f"{phf.idx}--{shape.name}--{phf.type}")
	shape.text = f"{phf.idx}--{shape.name}--{phf.type}"
	print("-------------------------------------------")
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	for shape in slide.placeholders:
		phf = shape.placeholder_format
		print(f"{phf.idx}--{shape.name}--{phf.type}")
		shape.text = f"{phf.idx}--{shape.name}--{phf.type}"

prs.save("哈哈.pptx")

# 段落的添加
import pptx
from pptx import Presentation

prs = Presentation("finall.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
name = slide.placeholders[14]
why = slide.placeholders[15]
name.text = "黄同学"
why.text = "学习太积极"
# --------------------------------------------------- #
prs1 = Presentation("finall.pptx")
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide1.shapes
title_shape = shapes.title # 这句代码可以改为title_shape = shapes.placeholders[0]
body_shape = shapes.placeholders[1]

title_shape.text = "这是一个标题"

tf = body_shape.text_frame
# 这句代码就是给body占位符添加内容！
tf.text = "带圆点的符号1"

p = tf.add_paragraph()
# 这个代码表示在原来的基础上，添加第一个段落！
p.text = "带圆点的符号2"

p = tf.add_paragraph()
# 这个代码表示在原来的基础上，添加第二个段落！
p.text = "带圆点的符号3"

prs.save("嘿嘿.pptx")

# 给段落设定层级关系
import pptx
from pptx import Presentation

prs = Presentation("finall.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
name = slide.placeholders[14]
why = slide.placeholders[15]
name.text = "黄同学"
why.text = "学习太积极"
# --------------------------------------------------- #
prs1 = Presentation("finall.pptx")
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide1.shapes
title_shape = shapes.title # 这句代码可以改为title_shape = shapes.placeholders[0]
body_shape = shapes.placeholders[1]

title_shape.text = "这是一个标题"

tf = body_shape.text_frame
tf.text = "带圆点的符号1"

p = tf.add_paragraph()
p.text = "带圆点的符号2"
# 原始内容的层级相当于是0，因此这个段落我设置为层级1，下面的段落设置为层级2
p.level = 1

p = tf.add_paragraph()
p.text = "带圆点的符号3"
p.level = 2

prs.save("嘻嘻.pptx")

# 添加一个文本框

# slide.shapes.add_textbox(left, top, width, height)
from pptx import Presentation
from pptx.util import Cm, Pt

prs = Presentation()
# 使用第一个版式
black_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是一段文本框里面的文字"

p = tf.add_paragraph()
p.text = "这是第二段文字，加粗，字号40"
p.font.bold = True
p.font.size = Pt(40)

prs.save("添加一个文本框0.pptx")

# 添加一个图片

# slide.shapes.add_picture(图片路径, 距离左边, 距离顶端, 宽度, 高度)
# 第一种展示：
from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = Cm(3)
pic = slide.shapes.add_picture("孙悟空.png", left, top)

prs.save("添加图片1.pptx")

# 第二种展示：
from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = Cm(3)
height = Cm(5.5)
pic = slide.shapes.add_picture("孙悟空.png", left, top, height=height)

prs.save("添加图片2.pptx")

# 添加表格

# shapes.add_table(rows, cols, left, top, width, height)

from pptx import Presentation
from pptx.util import Cm, Pt

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)
shapes = slide.shapes

rows, cols = 5, 3
left = top = Cm(5)
width = Cm(18)
height = Cm(3)

table = shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(2)
table.columns[2].width = Cm(2)
table.rows[0].height = Cm(2)

data = [
["姓名","性别","成绩"],
["张三","男",96],
["李四","女",87],
["王五","女",90],
["赵六","男",78]
]

for row in range(rows):
	for col in range(cols):
		table.cell(row,col).text = str(data[row][col])

prs.save("插入表格.pptx")

# 5、PPT文档内容样式批量调整
# 1）文本框位置的调整
# 上面我们已经知道怎么添加文本框，现在我们需要做的就是，怎么调整文本框的位置。
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是一段文本框里面的文字"
# ----------------------------------------- #
tf.margin_bottom = Cm(0.1) # 下边距
tf.margin_left = 0 # 下边距
# 一定要导入MSO_ANCHOR这个库
tf.vertical_anchor = MSO_ANCHOR.BOTTOM # 对齐文本方式：底端对齐
tf.word_wrap = True # 框中的文字自动换行

prs.save("文本框样式的调整.pptx")

# 2）文本框背景颜色调整
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是一段文本框里面的文字"
# -------------------------------------- #
tf.margin_bottom = Cm(0.1) # 下边距
tf.margin_left = 0 # 下边距
tf.vertical_anchor = MSO_ANCHOR.BOTTOM
tf.word_wrap = True # 框中的文字自动换行
# -------------------------------------- #
fill = text_box.fill
fill.solid()
# 使用之前一定要导入RGBColor这个库
fill.fore_color.rgb = RGBColor(247, 150, 70)

prs.save("文本框背景色的调整.pptx")

# 3）文本框边框样式调整
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是一段文本框里面的文字"
# -------------------------------------- #
tf.margin_bottom = Cm(0.1) # 下边距
tf.margin_left = 0 # 下边距
tf.vertical_anchor = MSO_ANCHOR.BOTTOM
tf.word_wrap = True # 框中的文字自动换行
# -------------------------------------- #
fill = text_box.fill
fill.solid()
# 使用之前一定要导入RGBColor这个库
fill.fore_color.rgb = RGBColor(247, 150, 70)
# -------------------------------------- #
line = text_box.line
line.color.rgb = RGBColor(255, 0, 0)
line.width = Cm(0.3)

prs.save("文本框边框样式调整.pptx")

# 4）段落对其调整
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
# ---------------------------- #
p = tf.add_paragraph()
p.text = "这是第二段文字"
p.alignment = PP_ALIGN.LEFT

prs.save("段落对其调整.pptx")

# 5）字体样式调整

# 代码如下：
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
# ---------------------------- #
p = tf.add_paragraph()
p.text = "这是第二段文字"
p.alignment = PP_ALIGN.LEFT
# ------------------------------------- #
p.font.bold = True
p.font.name = "宋体"
p.font.color.rgb = RGBColor(247, 150, 70)
p.font.size = Pt(30)

prs.save("字体样式调整.pptx")

